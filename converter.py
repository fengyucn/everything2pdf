"""
Everything to PDF Converter - 核心转换逻辑
支持图片、Office文档、PDF的转换和合并
多引擎支持：LibreOffice, Pandoc, Calibre, WeasyPrint, Python库
"""

import os
import io
import subprocess
import tempfile
import shutil
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any
from dataclasses import dataclass
from enum import Enum

import fitz  # PyMuPDF
from PIL import Image
import img2pdf

# 尝试导入Office文档处理库（用于后备方案）
try:
    from docx import Document
    from docx.shared import Inches
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

# 支持的文件格式
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.tif', '.webp'}
OFFICE_EXTENSIONS = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}
PDF_EXTENSION = '.pdf'

# 转换DPI设置
DEFAULT_DPI = 150


class ConversionEngine(Enum):
    """转换引擎枚举"""
    AUTO = "auto"  # 自动选择最佳引擎
    LIBREOFFICE = "libreoffice"  # LibreOffice（推荐，质量最好）
    PANDOC = "pandoc"  # Pandoc（轻量，学术文档友好）
    CALIBRE = "calibre"  # Calibre（电子书专家）
    WEASYPRINT = "weasyprint"  # WeasyPrint（HTML到PDF）
    PYTHON = "python"  # 纯Python实现（无外部依赖）


@dataclass
class EngineInfo:
    """引擎信息"""
    name: str
    display_name: str
    description: str
    available: bool
    priority: int  # 优先级，数字越小优先级越高
    supported_formats: List[str]


# 全局状态
_libreoffice_path = None
_libreoffice_checked = False
_engine_availability_cache = {}


def _run_libreoffice(lo_path: str, args: List[str], timeout: int = 120) -> subprocess.CompletedProcess:
    """运行LibreOffice命令，自动处理库路径"""
    env = os.environ.copy()
    lo_program_dir = os.path.join(os.path.dirname(os.path.realpath(lo_path)), '..', 'lib', 'libreoffice', 'program')
    if not os.path.isdir(lo_program_dir):
        lo_program_dir = '/usr/lib/libreoffice/program'
    if os.path.isdir(lo_program_dir):
        ld = env.get('LD_LIBRARY_PATH', '')
        env['LD_LIBRARY_PATH'] = lo_program_dir + (':' + ld if ld else '')
    return subprocess.run(args, capture_output=True, text=True, timeout=timeout, env=env)


def detect_libreoffice() -> Optional[str]:
    """检测系统中的LibreOffice安装路径（并验证可运行）"""
    # Linux常见路径
    linux_paths = [
        '/usr/bin/libreoffice',
        '/usr/bin/soffice',
        '/usr/local/bin/libreoffice',
        '/usr/local/bin/soffice',
        '/snap/bin/libreoffice',
    ]
    
    # Windows常见路径
    windows_paths = [
        r'C:\Program Files\LibreOffice\program\soffice.exe',
        r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
    ]
    
    # 根据操作系统选择路径列表
    if os.name == 'nt':
        paths = windows_paths
    else:
        paths = linux_paths
    
    candidates = []
    
    # 检查已知路径
    for path in paths:
        if os.path.isfile(path):
            candidates.append(path)
    
    # 尝试使用which/where命令查找
    if not candidates:
        try:
            if os.name == 'nt':
                result = subprocess.run(['where', 'soffice'], capture_output=True, text=True)
            else:
                result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True)
            
            if result.returncode == 0 and result.stdout.strip():
                candidates.append(result.stdout.strip().split('\n')[0])
        except Exception:
            pass
    
    # 验证候选路径是否能正常运行
    for path in candidates:
        try:
            result = _run_libreoffice(path, [path, '--headless', '--version'], timeout=10)
            if result.returncode == 0:
                return path
        except Exception:
            continue
    
    return None


def check_pandoc() -> bool:
    """检查Pandoc是否可用"""
    try:
        result = subprocess.run(['pandoc', '--version'], capture_output=True, timeout=5)
        return result.returncode == 0
    except Exception:
        return False


def check_calibre() -> bool:
    """检查Calibre是否可用"""
    try:
        result = subprocess.run(['ebook-convert', '--version'], capture_output=True, timeout=5)
        return result.returncode == 0
    except Exception:
        return False


def check_weasyprint() -> bool:
    """检查WeasyPrint是否可用"""
    try:
        import weasyprint
        return True
    except ImportError:
        return False


def get_engine_info() -> Dict[ConversionEngine, EngineInfo]:
    """获取所有引擎的信息"""
    global _engine_availability_cache
    
    if not _engine_availability_cache:
        # 检查LibreOffice（验证可运行）
        has_libreoffice, lo_path = get_libreoffice_status()
        
        _engine_availability_cache = {
            ConversionEngine.LIBREOFFICE: EngineInfo(
                name="libreoffice",
                display_name="LibreOffice",
                description="Office套件，转换质量最佳，支持所有Office格式",
                available=has_libreoffice,
                priority=1,
                supported_formats=['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
            ),
            ConversionEngine.PANDOC: EngineInfo(
                name="pandoc",
                display_name="Pandoc",
                description="轻量级文档转换器，学术文档友好",
                available=check_pandoc(),
                priority=2,
                supported_formats=['.doc', '.docx']
            ),
            ConversionEngine.CALIBRE: EngineInfo(
                name="calibre",
                display_name="Calibre",
                description="电子书专家，字体处理好",
                available=check_calibre(),
                priority=3,
                supported_formats=['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
            ),
            ConversionEngine.WEASYPRINT: EngineInfo(
                name="weasyprint",
                display_name="WeasyPrint",
                description="HTML到PDF转换器，CSS支持好",
                available=check_weasyprint(),
                priority=4,
                supported_formats=['.doc', '.docx']
            ),
            ConversionEngine.PYTHON: EngineInfo(
                name="python",
                display_name="Python库",
                description="纯Python实现，无需外部依赖",
                available=HAS_DOCX or HAS_OPENPYXL or HAS_PPTX,
                priority=5,
                supported_formats=['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
            )
        }
    
    return _engine_availability_cache


def get_available_engines() -> List[Dict[str, Any]]:
    """获取可用的引擎列表（用于前端显示）"""
    engines = get_engine_info()
    available_engines = []
    
    for engine, info in engines.items():
        if info.available:
            available_engines.append({
                'id': engine.value,
                'name': info.display_name,
                'description': info.description,
                'priority': info.priority,
                'supported_formats': info.supported_formats
            })
    
    # 按优先级排序
    available_engines.sort(key=lambda x: x['priority'])
    
    return available_engines


def get_file_type(filepath: str) -> str:
    """获取文件类型"""
    ext = Path(filepath).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return 'image'
    elif ext in OFFICE_EXTENSIONS:
        return 'office'
    elif ext == PDF_EXTENSION:
        return 'pdf'
    else:
        return 'unknown'


def convert_image_to_pdf_bytes(image_path: str) -> bytes:
    """将单个图片转换为PDF字节"""
    with open(image_path, 'rb') as f:
        img_bytes = f.read()
    
    # 检查图片格式，某些格式需要先转换
    try:
        img = Image.open(io.BytesIO(img_bytes))
        if img.mode in ('RGBA', 'P'):
            # 转换为RGB模式
            rgb_img = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'RGBA':
                rgb_img.paste(img, mask=img.split()[3])
            else:
                rgb_img.paste(img)
            
            buf = io.BytesIO()
            rgb_img.save(buf, format='JPEG', quality=95)
            img_bytes = buf.getvalue()
        elif img.format == 'GIF':
            # GIF转JPEG
            rgb_img = img.convert('RGB')
            buf = io.BytesIO()
            rgb_img.save(buf, format='JPEG', quality=95)
            img_bytes = buf.getvalue()
    except Exception:
        pass
    
    # 使用img2pdf转换
    pdf_bytes = img2pdf.convert(img_bytes)
    return pdf_bytes


def pdf_to_images(pdf_path: str, dpi: int = DEFAULT_DPI) -> List[bytes]:
    """将PDF转换为图片列表（PNG字节）"""
    images = []
    doc = fitz.open(pdf_path)
    
    zoom = dpi / 72  # 72是PDF默认DPI
    matrix = fitz.Matrix(zoom, zoom)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=matrix)
        img_bytes = pix.tobytes('png')
        images.append(img_bytes)
    
    doc.close()
    return images


# ==================== 多引擎转换函数 ====================

def convert_with_libreoffice(office_path: str, libreoffice_path: str, output_dir: str) -> Optional[str]:
    """使用LibreOffice转换Office文档为PDF"""
    try:
        cmd = [
            libreoffice_path,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            office_path
        ]
        
        result = _run_libreoffice(libreoffice_path, cmd)
        
        if result.returncode == 0:
            # 找到生成的PDF文件
            base_name = Path(office_path).stem
            pdf_path = os.path.join(output_dir, f'{base_name}.pdf')
            if os.path.exists(pdf_path):
                return pdf_path
    except Exception as e:
        print(f"LibreOffice转换失败: {e}")
    
    return None


def convert_with_pandoc(office_path: str, output_dir: str) -> Optional[str]:
    """使用Pandoc转换Office文档为PDF"""
    try:
        output_path = os.path.join(output_dir, f"{Path(office_path).stem}.pdf")
        
        cmd = [
            'pandoc',
            office_path,
            '-o', output_path,
            '--pdf-engine=xelatex',
            '-V', 'geometry:margin=1in',
            '-V', 'CJKmainfont=Noto Sans CJK SC'
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode == 0 and os.path.exists(output_path):
            return output_path
    except Exception as e:
        print(f"Pandoc转换失败: {e}")
    
    return None


def convert_with_calibre(office_path: str, output_dir: str) -> Optional[str]:
    """使用Calibre转换Office文档为PDF"""
    try:
        output_path = os.path.join(output_dir, f"{Path(office_path).stem}.pdf")
        
        cmd = [
            'ebook-convert',
            office_path,
            output_path,
            '--paper-size', 'a4',
            '--pdf-default-font-size', '12',
            '--pdf-mono-font-size', '10',
            '--pdf-page-numbers'
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        if result.returncode == 0 and os.path.exists(output_path):
            return output_path
    except Exception as e:
        print(f"Calibre转换失败: {e}")
    
    return None


def convert_with_weasyprint(office_path: str, output_dir: str) -> Optional[str]:
    """使用WeasyPrint转换Office文档为PDF"""
    try:
        from weasyprint import HTML, CSS
        
        output_path = os.path.join(output_dir, f"{Path(office_path).stem}.pdf")
        
        # 先用python-docx读取内容转HTML
        if office_path.endswith(('.doc', '.docx')):
            if not HAS_DOCX:
                return None
            
            doc = Document(office_path)
            html_content = convert_docx_to_html(doc)
        else:
            return None
        
        html = HTML(string=html_content)
        css = CSS(string='''
            @page { size: A4; margin: 1cm; }
            body { font-family: "Noto Sans CJK SC", "Microsoft YaHei", sans-serif; font-size: 12pt; }
            h1 { font-size: 18pt; color: #333; }
            h2 { font-size: 16pt; color: #444; }
            p { margin: 0.5em 0; line-height: 1.6; }
            table { border-collapse: collapse; width: 100%; margin: 1em 0; }
            td, th { border: 1px solid #ddd; padding: 8px; }
            th { background-color: #f2f2f2; }
        ''')
        html.write_pdf(output_path, stylesheets=[css])
        
        return output_path
    except Exception as e:
        print(f"WeasyPrint转换失败: {e}")
    
    return None


def convert_docx_to_html(doc) -> str:
    """将Word文档转换为HTML"""
    html_parts = ['<!DOCTYPE html><html><head><meta charset="utf-8"><title>Document</title></head><body>']
    
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            try:
                level = int(para.style.name.replace('Heading ', ''))
                html_parts.append(f'<h{level}>{para.text}</h{level}>')
            except:
                html_parts.append(f'<p>{para.text}</p>')
        else:
            if para.text.strip():
                html_parts.append(f'<p>{para.text}</p>')
    
    # 处理表格
    for table in doc.tables:
        html_parts.append('<table>')
        for row in table.rows:
            html_parts.append('<tr>')
            for cell in row.cells:
                html_parts.append(f'<td>{cell.text}</td>')
            html_parts.append('</tr>')
        html_parts.append('</table>')
    
    html_parts.append('</body></html>')
    return ''.join(html_parts)


def convert_docx_with_python(docx_path: str, output_dir: str) -> Optional[str]:
    """使用纯Python转换Word文档为PDF"""
    if not HAS_DOCX or not HAS_REPORTLAB:
        return None
    
    try:
        doc = Document(docx_path)
        output_path = os.path.join(output_dir, f'{Path(docx_path).stem}.pdf')
        
        # 尝试注册中文字体
        try:
            font_paths = [
                '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc',
                '/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf',
                '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
                'C:/Windows/Fonts/simhei.ttf',
                'C:/Windows/Fonts/msyh.ttc',
            ]
            for font_path in font_paths:
                if os.path.exists(font_path):
                    pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                    break
        except Exception:
            pass
        
        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        
        try:
            chinese_style = ParagraphStyle(
                'ChineseStyle',
                parent=styles['Normal'],
                fontName='ChineseFont',
                fontSize=12,
                leading=16,
            )
        except Exception:
            chinese_style = styles['Normal']
        
        story = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                try:
                    p = Paragraph(para.text, chinese_style)
                    story.append(p)
                    story.append(Spacer(1, 6))
                except Exception:
                    pass
        
        if story:
            pdf_doc.build(story)
            return output_path
        
    except Exception as e:
        print(f"Python Word转换失败: {e}")
    
    return None


def convert_xlsx_with_python(xlsx_path: str, output_dir: str) -> Optional[str]:
    """使用纯Python转换Excel文档为PDF"""
    if not HAS_OPENPYXL or not HAS_REPORTLAB:
        return None
    
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        output_path = os.path.join(output_dir, f'{Path(xlsx_path).stem}.pdf')
        
        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading2']))
            story.append(Spacer(1, 12))
            
            data = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                if any(row_data):
                    data.append(row_data)
            
            if data:
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                story.append(table)
            
            story.append(Spacer(1, 24))
        
        if story:
            pdf_doc.build(story)
            return output_path
        
    except Exception as e:
        print(f"Python Excel转换失败: {e}")
    
    return None


def convert_pptx_with_python(pptx_path: str, output_dir: str) -> Optional[str]:
    """使用纯Python转换PPT文档为PDF"""
    if not HAS_PPTX or not HAS_REPORTLAB:
        return None
    
    try:
        prs = Presentation(pptx_path)
        output_path = os.path.join(output_dir, f'{Path(pptx_path).stem}.pdf')
        
        pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            story.append(Paragraph(f"Slide {slide_num}", styles['Heading2']))
            story.append(Spacer(1, 12))
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    try:
                        p = Paragraph(shape.text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                    except Exception:
                        pass
            
            story.append(Spacer(1, 24))
        
        if story:
            pdf_doc.build(story)
            return output_path
        
    except Exception as e:
        print(f"Python PPT转换失败: {e}")
    
    return None


def convert_office_to_pdf(
    office_path: str,
    output_dir: str,
    engine: ConversionEngine = ConversionEngine.AUTO,
    libreoffice_path: Optional[str] = None
) -> Tuple[Optional[str], str]:
    """
    转换Office文档为PDF，支持多引擎选择
    
    Args:
        office_path: Office文档路径
        output_dir: 输出目录
        engine: 转换引擎，默认自动选择
        libreoffice_path: LibreOffice路径（可选）
    
    Returns:
        (PDF路径, 使用的引擎名称)
    """
    ext = Path(office_path).suffix.lower()
    engines_info = get_engine_info()
    
    def try_engine(eng: ConversionEngine) -> Optional[Tuple[str, str]]:
        """尝试用指定引擎转换，返回 (pdf路径, 引擎名) 或 None"""
        if eng == ConversionEngine.LIBREOFFICE and libreoffice_path:
            result = convert_with_libreoffice(office_path, libreoffice_path, output_dir)
            if result:
                return result, 'libreoffice'
        elif eng == ConversionEngine.PANDOC:
            result = convert_with_pandoc(office_path, output_dir)
            if result:
                return result, 'pandoc'
        elif eng == ConversionEngine.CALIBRE:
            result = convert_with_calibre(office_path, output_dir)
            if result:
                return result, 'calibre'
        elif eng == ConversionEngine.WEASYPRINT:
            result = convert_with_weasyprint(office_path, output_dir)
            if result:
                return result, 'weasyprint'
        elif eng == ConversionEngine.PYTHON:
            if ext in {'.doc', '.docx'}:
                result = convert_docx_with_python(office_path, output_dir)
            elif ext in {'.xls', '.xlsx'}:
                result = convert_xlsx_with_python(office_path, output_dir)
            elif ext in {'.ppt', '.pptx'}:
                result = convert_pptx_with_python(office_path, output_dir)
            else:
                result = None
            if result:
                return result, 'python'
        return None
    
    if engine == ConversionEngine.AUTO:
        # 自动模式：按优先级逐个尝试所有可用引擎
        sorted_engines = sorted(
            [(e, i) for e, i in engines_info.items()],
            key=lambda x: x[1].priority
        )
        for eng, info in sorted_engines:
            if info.available and ext in info.supported_formats:
                ret = try_engine(eng)
                if ret:
                    return ret
    else:
        # 指定引擎模式：先尝试指定的引擎
        ret = try_engine(engine)
        if ret:
            return ret
        # 指定引擎失败，回退到其他可用引擎
        sorted_engines = sorted(
            [(e, i) for e, i in engines_info.items()],
            key=lambda x: x[1].priority
        )
        for eng, info in sorted_engines:
            if eng != engine and info.available and ext in info.supported_formats:
                ret = try_engine(eng)
                if ret:
                    return ret
    
    return None, 'failed'


def images_to_pdf(image_bytes_list: List[bytes], output_path: str):
    """将图片字节列表合并为单个PDF"""
    doc = fitz.open()
    
    for img_bytes in image_bytes_list:
        # 从图片字节创建PDF页面
        img = fitz.open(stream=img_bytes, filetype='png')
        rect = img[0].rect
        
        # 创建新页面，大小与图片相同
        page = doc.new_page(width=rect.width, height=rect.height)
        page.insert_image(rect, stream=img_bytes)
        img.close()
    
    doc.save(output_path)
    doc.close()


def convert_files_to_pdf(
    file_paths: List[str],
    output_path: str,
    libreoffice_path: Optional[str] = None,
    engine: str = "auto",
    dpi: int = DEFAULT_DPI,
    progress_callback=None
) -> Tuple[bool, str]:
    """
    将多个文件转换并合并为单个PDF
    
    Args:
        file_paths: 文件路径列表（按顺序）
        output_path: 输出PDF路径
        libreoffice_path: LibreOffice路径（可选）
        engine: 转换引擎名称（auto/libreoffice/pandoc/calibre/weasyprint/python）
        dpi: 输出图片DPI
        progress_callback: 进度回调函数 callback(current, total, message)
    
    Returns:
        (成功标志, 消息)
    """
    if not file_paths:
        return False, "没有选择文件"
    
    # 解析引擎参数
    try:
        selected_engine = ConversionEngine(engine)
    except ValueError:
        selected_engine = ConversionEngine.AUTO
    
    all_images = []
    temp_dir = tempfile.mkdtemp()
    used_engines = set()
    
    try:
        total = len(file_paths)
        
        for idx, file_path in enumerate(file_paths):
            if progress_callback:
                progress_callback(idx, total, f"处理: {Path(file_path).name}")
            
            file_type = get_file_type(file_path)
            
            if file_type == 'image':
                # 图片：直接转为PDF字节，再转为图片
                try:
                    pdf_bytes = convert_image_to_pdf_bytes(file_path)
                    # 用PyMuPDF读取PDF字节并转为图片
                    temp_pdf = fitz.open(stream=pdf_bytes, filetype='pdf')
                    zoom = dpi / 72
                    matrix = fitz.Matrix(zoom, zoom)
                    for page in temp_pdf:
                        pix = page.get_pixmap(matrix=matrix)
                        all_images.append(pix.tobytes('png'))
                    temp_pdf.close()
                except Exception as e:
                    # 如果img2pdf失败，尝试直接用PIL
                    img = Image.open(file_path)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    buf = io.BytesIO()
                    img.save(buf, format='PNG')
                    all_images.append(buf.getvalue())
            
            elif file_type == 'office':
                # Office文档：使用多引擎转换
                pdf_path, used_engine = convert_office_to_pdf(
                    file_path, temp_dir, selected_engine, libreoffice_path
                )
                if pdf_path:
                    images = pdf_to_images(pdf_path, dpi)
                    all_images.extend(images)
                    used_engines.add(used_engine)
                else:
                    return False, f"无法转换文件: {Path(file_path).name}"
            
            elif file_type == 'pdf':
                # PDF：转为图片
                images = pdf_to_images(file_path, dpi)
                all_images.extend(images)
            
            else:
                return False, f"不支持的文件格式: {Path(file_path).name}"
        
        if progress_callback:
            progress_callback(total, total, "正在生成PDF...")
        
        # 合并所有图片为PDF
        if all_images:
            images_to_pdf(all_images, output_path)
            engines_str = ", ".join(sorted(used_engines)) if used_engines else "直接处理"
            return True, f"成功转换 {len(file_paths)} 个文件 (引擎: {engines_str})"
        else:
            return False, "没有可转换的内容"
    
    except Exception as e:
        return False, f"转换失败: {str(e)}"
    
    finally:
        # 清理临时目录
        try:
            shutil.rmtree(temp_dir)
        except Exception:
            pass


# ==================== 状态查询函数 ====================

def get_libreoffice_status() -> Tuple[bool, Optional[str]]:
    """获取LibreOffice状态"""
    global _libreoffice_path, _libreoffice_checked
    
    if not _libreoffice_checked:
        _libreoffice_path = detect_libreoffice()
        _libreoffice_checked = True
    
    return _libreoffice_path is not None, _libreoffice_path


def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名列表"""
    extensions = list(IMAGE_EXTENSIONS) + list(OFFICE_EXTENSIONS) + [PDF_EXTENSION]
    return sorted(extensions)


def get_conversion_engines_status() -> Dict[str, Any]:
    """获取所有转换引擎的状态信息"""
    engines = get_engine_info()
    
    return {
        'engines': [
            {
                'id': engine.value,
                'name': info.display_name,
                'description': info.description,
                'available': info.available,
                'priority': info.priority,
                'supported_formats': info.supported_formats
            }
            for engine, info in engines.items()
        ],
        'recommended': next(
            (engine.value for engine, info in engines.items() if info.available),
            'python'
        )
    }
